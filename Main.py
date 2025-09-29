from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Cm
import argparse
import util



# Create parser
parser = argparse.ArgumentParser(description="Demo of argparse")
# Add arguments
parser.add_argument("--filename", type=str, required=True)
args = parser.parse_args()
file_name = args.filename
file_name_suffix = file_name.split(".")[0]
## Load pptx file 
prs = Presentation(file_name)

for slide_num, slide in enumerate(prs.slides, start=1):
    slide_height = prs.slide_height
    slide_width = prs.slide_width / util.EMU_PER_CM
    print(f"Processing --- Slide {slide_num} ---")
    with open(f"output_{file_name_suffix}_{slide_num}.tex", "w", encoding="utf-8") as f:
        # f.write("\\resizebox{0.9\linewidth}{!}{%\n")
        f.write("\\begin{tikzpicture}[scale=\\textwidth/")
        f.write(f"{slide_width:.2f}cm")
        f.write(", every node/.style={transform shape=false}]\n")
        for shape in slide.shapes:
            tikz_cmd = util.shape_to_tikz(shape, slide_height)
            f.write(tikz_cmd)
        f.write("\\end{tikzpicture}\n")
    f.close()
    print(f"--- Generated file : output_{file_name_suffix}_{slide_num}.tex ---")
print(" Done !")

       